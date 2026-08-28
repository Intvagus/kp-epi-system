"""One-off asset-generation script: traces src/dashboard/kp_districts.geojson
from the user-provided reference map (KP_MAP_1.pptx), replacing the earlier
geoBoundaries.org set. Not part of the runtime pipeline/dashboard -- run this
only if the reference map changes and the geojson needs regenerating.

The pptx's district labels are themselves outlined-to-vector text (a Google
Slides export artifact, not real PowerPoint text runs), so shape identity was
confirmed by rendering the pptx's own vector geometry directly (district
polygons + label outlines overlaid) and reading it visually -- see this
project's session notes for the full verification, including the one
unconfirmed assumption (SW Wazir Belt vs SW Mehsud Belt, flagged in
config.py's DISTRICT_TO_BOUNDARY). SHAPE_TO_DISTRICT below is the result.

Requires: pip install python-pptx shapely (shapely is NOT a runtime
dependency of this project -- only needed to re-run this script).

Usage: python scripts/extract_kp_districts_geojson.py path/to/KP_MAP_1.pptx
Writes: src/dashboard/kp_districts.geojson
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from shapely.geometry import Polygon as ShapelyPolygon

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = PROJECT_ROOT / "src" / "dashboard" / "kp_districts.geojson"

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
      "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

# Verified by direct visual inspection of KP_MAP_1.pptx. Shape IDs are the
# slide's real PowerPoint shape IDs. Names match src/pipeline/config.py's
# DISTRICT_TO_BOUNDARY keys exactly, plus "Tor Ghar" (a real district with
# no coverage data in the files received -- see CLAUDE.md).
SHAPE_TO_DISTRICT = {
    5: "Abbottabad", 6: "Bajaur", 8: "Bannu", 9: "Battagram", 10: "Buner",
    11: "Charsadda", 13: "Chitral Lower", 14: "Chitral Upper", 15: "D.I. Khan",
    17: "Hangu", 20: "Haripur", 23: "Karak", 25: "Khyber", 26: "Kohat",
    27: "Kohistan Lower", 28: "Kohistan Upper", 29: "Kolai Palas Kohistan",
    30: "Kurram Lower and Central", 31: "Lakki Marwat", 33: "Dir Lower",
    34: "SW Mehsud Belt", 35: "Malakand", 36: "Mansehra", 38: "Mardan",
    40: "Mohmand", 43: "North Waziristan", 44: "Nowshera", 46: "Orakzai",
    47: "Peshawar", 49: "Shangla", 50: "Swabi", 52: "Swat", 54: "Tank",
    55: "Tor Ghar", 56: "Dir Upper", 57: "Kurram Upper", 58: "SW Wazir Belt",
}


def _get_xfrm(sp_el):
    xfrm = sp_el.find("./p:spPr/a:xfrm", NS)
    if xfrm is None:
        xfrm = sp_el.find("./p:grpSpPr/a:xfrm", NS)
    if xfrm is None:
        return None
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    choff = xfrm.find("a:chOff", NS)
    chext = xfrm.find("a:chExt", NS)
    d = {"off": (int(off.get("x")), int(off.get("y"))), "ext": (int(ext.get("cx")), int(ext.get("cy")))}
    if choff is not None and chext is not None:
        d["choff"] = (int(choff.get("x")), int(choff.get("y")))
        d["chext"] = (int(chext.get("cx")), int(chext.get("cy")))
    return d


def _cubic_bezier_flatten(p0, p1, p2, p3, n=6):
    pts = []
    for i in range(1, n + 1):
        t = i / n
        mt = 1 - t
        x = (mt ** 3) * p0[0] + 3 * (mt ** 2) * t * p1[0] + 3 * mt * (t ** 2) * p2[0] + (t ** 3) * p3[0]
        y = (mt ** 3) * p0[1] + 3 * (mt ** 2) * t * p1[1] + 3 * mt * (t ** 2) * p2[1] + (t ** 3) * p3[1]
        pts.append((x, y))
    return pts


def _path_rings(sp_el):
    """(w, h, ring_points) per <a:path> -- dense polyline vertices in
    path-local coordinate space, cubic beziers flattened."""
    rings = []
    for path in sp_el.findall("./p:spPr/a:custGeom/a:pathLst/a:path", NS):
        w = int(path.get("w", 1))
        h = int(path.get("h", 1))
        pts, cur = [], None
        for child in path:
            tag = child.tag.split("}")[-1]
            if tag in ("moveTo", "lnTo"):
                pt = child.find("a:pt", NS)
                cur = (int(pt.get("x")), int(pt.get("y")))
                pts.append(cur)
            elif tag == "cubicBezTo":
                cps = [(int(p.get("x")), int(p.get("y"))) for p in child.findall("a:pt", NS)]
                pts.extend(_cubic_bezier_flatten(cur, cps[0], cps[1], cps[2]))
                cur = cps[2]
        if pts:
            rings.append((w, h, pts))
    return rings


def _transform(x, y, off, ext, w, h):
    return off[0] + (x / w) * ext[0], off[1] + (y / h) * ext[1]


def _group_transform(x, y, g):
    fx = (x - g["choff"][0]) / g["chext"][0]
    fy = (y - g["choff"][1]) / g["chext"][1]
    return g["off"][0] + fx * g["ext"][0], g["off"][1] + fy * g["ext"][1]


def extract(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))
    slide_el = prs.slides[0].shapes._spTree
    features = []

    def walk(el, group_xfrm, depth):
        for sp in el:
            tag = sp.tag.split("}")[-1]
            if tag == "grpSp":
                walk(sp, _get_xfrm(sp), depth + 1)
            elif tag == "sp":
                sid = int(sp.find(".//p:cNvPr", NS).get("id"))
                if sid not in SHAPE_TO_DISTRICT:
                    continue
                xfrm = _get_xfrm(sp)
                polygon_rings = []
                for w, h, pts in _path_rings(sp):
                    ring = []
                    for (x, y) in pts:
                        sx, sy = _transform(x, y, xfrm["off"], xfrm["ext"], w, h)
                        if group_xfrm is not None:
                            sx, sy = _group_transform(sx, sy, group_xfrm)
                        # inches, y flipped so north is up
                        ring.append([round(sx / 914400, 5), round(-sy / 914400, 5)])
                    if ring and ring[0] != ring[-1]:
                        ring.append(ring[0])
                    polygon_rings.append(ring)
                # Douglas-Peucker simplify -- these paths came from Google
                # Slides' own dense bezier-outlined export, far denser than a
                # ~640x520 SVG choropleth needs; cuts embedded dashboard.html
                # size substantially with no visible loss at map scale.
                try:
                    shp = ShapelyPolygon(polygon_rings[0], polygon_rings[1:])
                    simplified = shp.simplify(0.008, preserve_topology=True)
                    if simplified.geom_type == "Polygon" and not simplified.is_empty:
                        ext_ring = [[round(x, 4), round(y, 4)] for x, y in simplified.exterior.coords]
                        int_rings = [[[round(x, 4), round(y, 4)] for x, y in r.coords] for r in simplified.interiors]
                        polygon_rings = [ext_ring] + int_rings
                except Exception as e:
                    print(f"  simplify failed for {SHAPE_TO_DISTRICT[sid]}: {e}")
                features.append({
                    "type": "Feature",
                    "properties": {"shapeName": SHAPE_TO_DISTRICT[sid]},
                    "geometry": {"type": "Polygon", "coordinates": polygon_rings},
                })

    walk(slide_el, None, 0)

    found = {f["properties"]["shapeName"] for f in features}
    missing = set(SHAPE_TO_DISTRICT.values()) - found
    if missing:
        raise SystemExit(f"Missing shapes for: {sorted(missing)} -- check SHAPE_TO_DISTRICT against the source pptx")

    return {"type": "FeatureCollection", "features": features}


if __name__ == "__main__":
    if len(sys.argv) != 2:
        raise SystemExit("Usage: python scripts/extract_kp_districts_geojson.py path/to/KP_MAP_1.pptx")
    gj = extract(Path(sys.argv[1]))
    with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
        json.dump(gj, f)
    print(f"Wrote {OUTPUT_PATH} ({len(gj['features'])} features, {OUTPUT_PATH.stat().st_size / 1024:.0f} KB)")
