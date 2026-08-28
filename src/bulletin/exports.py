"""Excel annex (full tables, for recipients to re-cut the numbers) and a
PowerPoint deck (key charts/numbers for review meetings), both generated
from the SAME bulletin dict the PDF renders from -- data/processed/vpd_summary.json's
'bulletin' section. Word export is intentionally not implemented: the build
brief asks for it only on request, warning the layout won't survive Word.
"""
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter

PROJECT_ROOT = Path(__file__).resolve().parents[2]


NAVY = "16326B"
HEADER_FILL = PatternFill(start_color=NAVY, end_color=NAVY, fill_type="solid")
HEADER_FONT = Font(color="FFFFFF", bold=True)
TITLE_FONT = Font(bold=True, size=13, color=NAVY)


def _write_table(ws, start_row, headers, rows, title=None):
    r = start_row
    if title:
        ws.cell(row=r, column=1, value=title).font = Font(bold=True, size=11, color=NAVY)
        r += 1
    for c, h in enumerate(headers, start=1):
        cell = ws.cell(row=r, column=c, value=h)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
    r += 1
    for row in rows:
        for c, v in enumerate(row, start=1):
            ws.cell(row=r, column=c, value=v)
        r += 1
    for c in range(1, len(headers) + 1):
        ws.column_dimensions[get_column_letter(c)].width = max(16, len(str(headers[c - 1])) + 4)
    return r + 1  # next free row


def build_excel_annex(bulletin: dict, output_dir: Path) -> Path:
    wb = Workbook()

    ws = wb.active
    ws.title = "Summary"
    ws.cell(row=1, column=1, value=f"KP EPI Surveillance Bulletin — {bulletin['week_label']}").font = TITLE_FONT
    ws.cell(row=2, column=1, value=f"Cumulative: {bulletin['cumulative_label']}")
    r = _write_table(ws, 4, ["Indicator", "Value"], [
        ["Suspected measles-rubella cases (this week)", bulletin["msl"]["suspected_this_week"]],
        ["Measles confirmed (this week)", bulletin["msl"]["measles_confirmed_this_week"]],
        ["Rubella confirmed (this week)", bulletin["msl"]["rubella_confirmed_this_week"]],
        ["% zero-dose among eligible suspected (this week)", bulletin["msl"]["pct_zero_dose_eligible_this_week"]],
        ["Total suspected (cumulative)", bulletin["msl"]["cumulative"]["suspected"]],
        ["Total measles confirmed (cumulative)", bulletin["msl"]["cumulative"]["measles_confirmed"]],
        ["Total rubella confirmed (cumulative)", bulletin["msl"]["cumulative"]["rubella_confirmed"]],
        ["Sample collection rate % (cumulative)", bulletin["msl"]["cumulative"]["sample_collection_rate_pct"]],
        ["Diphtheria cases (this week)", bulletin["diphtheria"]["this_week"]["case_count"]],
        ["Diphtheria cases (cumulative)", bulletin["diphtheria"]["cumulative"]["case_count"]],
        ["Diphtheria districts affected (cumulative)", bulletin["diphtheria"]["cumulative"]["districts_affected"]],
        ["Diphtheria deaths (cumulative)", bulletin["diphtheria"]["cumulative"]["deaths"]],
        ["Pertussis cases (cumulative)", bulletin["pertussis"]["grand_total_cumulative"]],
        ["NNT cases (this week)", bulletin["nnt"]["this_week_case_count"]],
        ["NNT cases (cumulative)", bulletin["nnt"]["cumulative_case_count"]],
    ], title="Key indicators")
    note = ("Note: Measles/rubella incidence-rate indicators are not included -- no population "
            "denominator is available in the source data (see CLAUDE.md 'Confirmed VPD decisions').")
    ws.cell(row=r, column=1, value=note).font = Font(italic=True, size=9, color="666666")
    ws.column_dimensions["A"].width = 46

    ws2 = wb.create_sheet("MSL Classification")
    _write_table(ws2, 1, ["Final classification", "Count (this week)"],
                 sorted(bulletin["msl"]["classification_breakdown_this_week"].items(), key=lambda x: -x[1]),
                 title=f"Final classification — {bulletin['week_label']}")

    ws3 = wb.create_sheet("MSL Dose Status")
    dose_rows = []
    susp, conf = bulletin["msl"]["dose_status_suspected_this_week"], bulletin["msl"]["dose_status_confirmed_this_week"]
    for k in ["Zero dose", "1 dose", "2 doses", "Unknown"]:
        dose_rows.append([k, susp.get(k, 0), conf.get(k, 0)])
    _write_table(ws3, 1, ["Dose status", "Suspected (eligible, age 9m+)", "Confirmed (eligible, age 9m+)"],
                 dose_rows, title=f"MCV dose status — {bulletin['week_label']}")

    ws4 = wb.create_sheet("MSL Age Distribution")
    age_order = ["0-8m", "9-23m", "24-59m", "60m+"]
    age_dist = bulletin["msl"]["age_distribution_this_week"]
    _write_table(ws4, 1, ["Age group", "Suspected cases"],
                 [[k, age_dist.get(k, 0)] for k in age_order], title=f"Age distribution — {bulletin['week_label']}")

    ws5 = wb.create_sheet("MSL District Breakdown")
    d_rows = sorted(bulletin["msl"]["district_breakdown_this_week"], key=lambda r: -r["suspected"])
    _write_table(ws5, 1, ["District", "Suspected", "Measles confirmed", "Rubella confirmed"],
                 [[r["district"], r["suspected"], r["measles_confirmed"], r["rubella_confirmed"]] for r in d_rows],
                 title=f"District-wise — {bulletin['week_label']}")

    ws6 = wb.create_sheet("Diphtheria")
    diph = bulletin["diphtheria"]
    r = _write_table(ws6, 1, ["Indicator", "This week", "Cumulative"], [
        ["Cases", diph["this_week"]["case_count"], diph["cumulative"]["case_count"]],
        ["Districts affected", diph["this_week"]["districts_affected"], diph["cumulative"]["districts_affected"]],
        ["Lab confirmed", diph["this_week"]["lab_confirmed_count"], diph["cumulative"]["lab_confirmed_count"]],
        ["% no DPT history", diph["this_week"]["pct_no_dpt_history"], diph["cumulative"]["pct_no_dpt_history"]],
        ["% aged 5+", diph["this_week"]["pct_aged_5_plus"], diph["cumulative"]["pct_aged_5_plus"]],
        ["Deaths", diph["this_week"]["deaths"], diph["cumulative"]["deaths"]],
    ], title="Diphtheria summary")
    age_bands = sorted(diph["age_band_counts_cumulative"].items(),
                        key=lambda x: -1 if x[0] == "<1y" else int(x[0].split("-")[0]))
    _write_table(ws6, r, ["Age band", "Cases (cumulative)"], age_bands, title=f"Age bands — {bulletin['cumulative_label']}")

    ws7 = wb.create_sheet("Pertussis")
    p_rows = sorted(bulletin["pertussis"]["district_breakdown_cumulative"], key=lambda r: -r["case_count"])
    _write_table(ws7, 1, ["District", "Cases (cumulative)"],
                 [[r["district"], r["case_count"]] for r in p_rows],
                 title=f"Pertussis by district — {bulletin['cumulative_label']}")

    ws8 = wb.create_sheet("NNT")
    n_rows = sorted(bulletin["nnt"]["cumulative_district_breakdown"], key=lambda r: r["district"])
    _write_table(ws8, 1, ["District", "Cases (cumulative)"],
                 [[r["district"], r["case_count"]] for r in n_rows],
                 title=f"NNT by district — {bulletin['cumulative_label']}")

    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"Bulletin_Week_{bulletin['epi_week']}_{bulletin['year']}_annex.xlsx"
    wb.save(out_path)
    return out_path


def build_pptx(bulletin: dict, output_dir: Path) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    NAVY_RGB = RGBColor(0x16, 0x32, 0x6B)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def add_title(slide, text, subtitle=None):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(30)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY_RGB
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    def add_bullets(slide, items, top=1.5, left=0.6, width=12.0, size=18):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + item
            p.font.size = Pt(size)
            p.space_after = Pt(8)

    def add_table(slide, headers, rows, top=1.6, left=0.6, width=12.0, height=5.0):
        n_rows, n_cols = len(rows) + 1, len(headers)
        shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
        table = shape.table
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = str(h)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY_RGB
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for r, row in enumerate(rows, start=1):
            for c, v in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(v)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
        return table

    # Slide 1: Title
    s = add_slide()
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.2))
    tf = box.text_frame
    tf.text = "Surveillance Weekly Bulletin"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY_RGB
    p = tf.add_paragraph()
    p.text = "Vaccine Preventable Diseases (VPD)"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xB0, 0x8A, 0x00)
    p2 = tf.add_paragraph()
    p2.text = bulletin["week_label"] + " — Directorate General Health Services, KP EPI Section"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Slide 2: Measles-Rubella highlights
    s = add_slide()
    msl = bulletin["msl"]
    add_title(s, "Measles-Rubella Highlights", bulletin["week_label"])
    wow = msl["week_over_week"]
    bullets = [
        f"{msl['suspected_this_week']:,} suspected cases reported this week.",
        f"{msl['measles_confirmed_this_week']} measles lab confirmed, {msl['rubella_confirmed_this_week']} rubella confirmed.",
        (f"{msl['pct_zero_dose_eligible_this_week']:.0f}% of eligible suspected cases were zero-dose for MCV."
         if msl["pct_zero_dose_eligible_this_week"] is not None else "No eligible suspected cases this week."),
        f"Cumulative ({bulletin['cumulative_label']}): {msl['cumulative']['suspected']:,} suspected, "
        f"{msl['cumulative']['measles_confirmed']:,} measles confirmed, {msl['cumulative']['rubella_confirmed']:,} rubella confirmed.",
        f"Sample collection rate: {msl['cumulative']['sample_collection_rate_pct']:.0f}%.",
    ]
    if wow:
        bullets.append(f"Week-over-week: {wow['this_week_count']} this week vs {wow['prior_week_count']} in Week {wow['prior_week']}.")
    add_bullets(s, bullets)

    # Slide 3: Diphtheria
    s = add_slide()
    diph = bulletin["diphtheria"]
    add_title(s, "Diphtheria Situation", bulletin["cumulative_label"])
    this_week_n = diph["this_week"]["case_count"]
    deaths_n = diph["cumulative"]["deaths"]
    add_bullets(s, [
        f"{this_week_n} {'case' if this_week_n == 1 else 'cases'} reported this week.",
        f"{diph['cumulative']['case_count']} cases reported from {diph['cumulative']['districts_affected']} districts, cumulative.",
        f"{diph['cumulative']['lab_confirmed_count']} lab confirmed.",
        f"{diph['cumulative']['pct_no_dpt_history']:.0f}% of cases have no recorded DPT history.",
        f"{diph['cumulative']['pct_aged_5_plus']:.0f}% of cases are aged 5 years and above.",
        f"{deaths_n} diphtheria-related {'death' if deaths_n == 1 else 'deaths'} reported to date.",
    ])

    # Slide 4: Pertussis & NNT
    s = add_slide()
    add_title(s, "Pertussis & Neonatal Tetanus (NNT)", bulletin["cumulative_label"])
    pert_top = sorted(bulletin["pertussis"]["district_breakdown_cumulative"], key=lambda r: -r["case_count"])[:8]
    add_table(s, ["District", "Pertussis cases"], [[r["district"], r["case_count"]] for r in pert_top],
              top=1.6, left=0.6, width=5.8, height=4.5)
    nnt_top = sorted(bulletin["nnt"]["cumulative_district_breakdown"], key=lambda r: -r["case_count"])[:8]
    add_table(s, ["District", "NNT cases"], [[r["district"], r["case_count"]] for r in nnt_top],
              top=1.6, left=6.9, width=5.8, height=4.5)

    # Slide 5: AFP & Key Messages
    s = add_slide()
    add_title(s, "AFP Surveillance & Key Messages")
    add_bullets(s, [bulletin["afp"]["message"]] + ["Key messages:"] + bulletin["key_messages"], size=16)

    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"Bulletin_Week_{bulletin['epi_week']}_{bulletin['year']}.pptx"
    prs.save(out_path)
    return out_path
